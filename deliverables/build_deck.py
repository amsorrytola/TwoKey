"""Interlock business proposal deck.

Design brief, from current practice on minimal presentation design:
  one idea per slide, two typefaces, a tight palette, generous whitespace,
  restraint over decoration, and consistency from first slide to last.

Palette is Accenture's own: #A100FF purple, black, white. Nothing else except
the five lane-state colours, where colour carries data rather than decoration.
Type is IBM Plex Sans and JetBrains Mono, matching the product interface.
Titles are large, body text is never below 14pt, and no slide carries more than
one argument.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── palette ──────────────────────────────────────────────────────────────────
PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DEEP   = RGBColor(0x7A, 0x00, 0xC2)
BLACK  = RGBColor(0x08, 0x08, 0x0A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x16, 0x17, 0x1A)
INK2   = RGBColor(0x55, 0x57, 0x5E)
INK3   = RGBColor(0x9A, 0x9C, 0xA4)
RULE   = RGBColor(0xE4, 0xE4, 0xE8)
WASH   = RGBColor(0xFA, 0xF6, 0xFF)
GREEN  = RGBColor(0x06, 0x7A, 0x55)
AMBER  = RGBColor(0x9A, 0x62, 0x00)
RED    = RGBColor(0xC1, 0x1B, 0x1B)
BLUE   = RGBColor(0x0B, 0x5F, 0xCC)
DIM    = RGBColor(0x8E, 0x91, 0x9A)   # readable on black; 0x6A was too dark

# Typeface choice is an engineering decision, not an aesthetic one. A deck is opened
# on machines we do not control, and PowerPoint silently falls back to a serif when the
# named face is missing, which is what happened with IBM Plex Sans. Arial is present on
# every Mac and every Windows install, renders identically on both, and is the Swiss
# grotesque this layout was designed around in the first place. Courier New is reserved
# for genuine code strings, where monospacing carries meaning.
SANS, MONO = "Arial", "Courier New"
W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)                       # generous margin, the whitespace budget
CW = W - 2 * M                         # content width


# ── primitives ───────────────────────────────────────────────────────────────
def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, 0, 0, W, H, fill=bg)
    return s


def box(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    return sh


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.0):
    tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, r in enumerate(runs):
        t, size, bold, colour, font, after = (list(r) + [None] * 6)[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        if after:
            p.space_after = Pt(after)
        run = p.add_run()
        run.text = t
        run.font.size = Pt(size or 16)
        run.font.bold = bool(bold)
        run.font.color.rgb = colour or INK
        run.font.name = font or SANS
    return tb


def title(s, text, sub=None, dark=False, y=Inches(0.72)):
    """One line, large. The slide's single idea."""
    ink = WHITE if dark else INK
    ink_s = INK3 if dark else INK2
    box(s, M, y, Inches(0.62), Emu(28575), fill=PURPLE)
    txt(s, M, y + Inches(0.28), CW, Inches(1.1),
        [(text, 32, True, ink, SANS, 0)], ls=1.08)
    yy = y + Inches(1.12)
    if sub:
        txt(s, M, yy, Inches(10.4), Inches(0.7),
            [(sub, 16, False, ink_s, SANS, 0)], ls=1.3)
        yy += Inches(0.72)
    return yy + Inches(0.34)


def num(s, n):
    txt(s, W - M - Inches(0.7), Inches(0.62), Inches(0.7), Inches(0.4),
        [(f"{n:02d}", 13, True, RULE, SANS, 0)], align=PP_ALIGN.RIGHT)


def rows(s, x, y, w, data, widths, rh=Inches(0.62), size=14, headers=None):
    """Rule-separated rows. No fills, no zebra striping, no borders."""
    tot = sum(widths)
    xs, acc = [], x
    for cwv in widths:
        xs.append(acc); acc += int(w * cwv / tot)
    yy = y
    if headers:
        for i, h in enumerate(headers):
            txt(s, xs[i], yy, int(w * widths[i] / tot) - Inches(0.18), Inches(0.3),
                [(h.upper(), 10.5, True, INK3, SANS, 0)])
        yy += Inches(0.34)
        box(s, x, yy, w, Emu(9525), fill=INK)
        yy += Inches(0.16)
    for r in data:
        for i, cell in enumerate(r):
            t, colour, mono, bold = (list(cell) + [None] * 4)[:4]
            txt(s, xs[i], yy, int(w * widths[i] / tot) - Inches(0.18), rh,
                [(str(t), size, bold if bold is not None else False,
                  colour or INK, MONO if mono else SANS, 0)], ls=1.2)
        yy += rh
        box(s, x, yy - Inches(0.13), w, Emu(9525), fill=RULE)
    return yy


def stat(s, x, y, w, value, label, colour=None, vsize=54):
    """A single number, given room. No card, no border, no fill."""
    txt(s, x, y, w, Inches(0.9), [(value, vsize, True, colour or INK, SANS, 0)], ls=0.95)
    txt(s, x, y + Inches(0.86), w, Inches(0.8),
        [(label, 13, False, INK2, SANS, 0)], ls=1.25)


def quad(s, y, items, vsize=48):
    """Four statistics across, separated by whitespace rather than boxes."""
    n = len(items)
    gap = Inches(0.4)
    cw = int((CW - gap * (n - 1)) / n)
    for i, (v, lab, colour) in enumerate(items):
        x = M + i * (cw + gap)
        box(s, x, y, Inches(0.5), Emu(28575), fill=PURPLE)
        stat(s, x, y + Inches(0.3), cw, v, lab, colour, vsize)


def foot(s, text, dark=False):
    """Anchored above the bottom edge, with room for two wrapped lines."""
    txt(s, M, H - Inches(0.86), CW, Inches(0.56),
        [(text, 13, False, INK3 if dark else INK2, SANS, 0)], ls=1.25)


# ══════════════════════════════════════════════════════════════════ slides ═══
def build():
    prs = deck()

    # ── 1 cover ──────────────────────────────────────────────────────────────
    s = slide(prs, BLACK)
    box(s, 0, 0, W, Inches(0.1), fill=PURPLE)
    txt(s, M, Inches(2.45), Inches(1.2), Inches(0.9),
        [(">", 62, True, PURPLE, MONO, 0)], ls=0.9)
    txt(s, M, Inches(3.35), Inches(9), Inches(1.1),
        [("Interlock", 62, True, WHITE, SANS, 0)], ls=0.95)
    txt(s, M, Inches(4.55), Inches(8.6), Inches(0.9),
        [("Runtime assurance for agentic AI.", 21, False, INK3, SANS, 0)], ls=1.3)
    box(s, M, Inches(5.85), Inches(2.2), Emu(9525), fill=RGBColor(0x33, 0x33, 0x38))
    txt(s, M, Inches(6.1), Inches(9), Inches(0.9),
        [("Team TwoKey", 14, True, WHITE, SANS, 4),
         ("Accenture Innovation Challenge 2026  ·  Round 2  ·  ControlPlane.ai",
          12.5, False, DIM, SANS, 0)], ls=1.3)

    # ── 2 team ──────────────────────────────────────────────────────────────
    s = slide(prs)
    y = title(s, "Team details")
    box(s, M, y, CW, Inches(0.52), fill=WASH)
    txt(s, M + Inches(0.3), y + Inches(0.14), Inches(3), Inches(0.3),
        [("TEAM NAME", 11.5, True, INK2, SANS, 0)])
    txt(s, M + Inches(3.4), y + Inches(0.09), Inches(6), Inches(0.4),
        [("TwoKey", 18, True, PURPLE, SANS, 0)])

    y2 = y + Inches(1.0)
    photo = Inches(1.65)                       # both portraits at one size
    for i, (name, role, img) in enumerate(
            [("Riddhi Sidana", "Team Leader", "assets/riddhi.jpg"),
             ("Mohammed Talha Ansari", "Member", "assets/talha.jpg")]):
        x = M + i * (Inches(5.5) + Inches(0.6))
        box(s, x, y2, photo + Inches(0.09), photo + Inches(0.09), fill=PURPLE)
        s.shapes.add_picture(img, int(x + Inches(0.045)), int(y2 + Inches(0.045)),
                             int(photo), int(photo))
        tx = x + photo + Inches(0.42)
        txt(s, tx, y2 + Inches(0.02), Inches(3.4), Inches(0.5),
            [(name, 19, True, INK, SANS, 0)], ls=1.1)
        txt(s, tx, y2 + Inches(0.52), Inches(3.4), Inches(0.3),
            [(role, 12.5, True, PURPLE, SANS, 0)])
        for j, (k, v) in enumerate([("College", "IIT Roorkee"),
                                    ("Stream", "BS-MS Economics"),
                                    ("Year", "Fourth year")]):
            ry = y2 + Inches(0.95 + j * 0.3)
            txt(s, tx, ry, Inches(1.0), Inches(0.28),
                [(k, 11.5, True, INK2, SANS, 0)])
            txt(s, tx + Inches(1.05), ry, Inches(2.4), Inches(0.28),
                [(v, 11.5, False, INK, SANS, 0)])

    y3 = y2 + Inches(2.25)
    box(s, M, y3, CW, Emu(9525), fill=RULE)
    txt(s, M, y3 + Inches(0.25), CW, Inches(0.35),
        [("WHAT WE BUILT FOR ROUND 2", 11.5, True, INK3, SANS, 0)])
    facts = [("Working prototype", "Six-check mesh, five lanes, signed ledger, four-screen console"),
             ("Real models", "Four models, two vendors, two independent credentials"),
             ("Measured, not claimed", "Precision, recall, latency and cost from a reproducible harness")]
    cwf = int((CW - Inches(0.8)) / 3)
    for i, (h, b) in enumerate(facts):
        x = M + i * (cwf + Inches(0.4))
        txt(s, x, y3 + Inches(0.68), cwf, Inches(0.3),
            [(h, 14, True, INK, SANS, 0)])
        txt(s, x, y3 + Inches(1.0), cwf, Inches(0.6),
            [(b, 12, False, INK2, SANS, 0)], ls=1.3)

    # ── 3 the statement ──────────────────────────────────────────────────────
    s = slide(prs, BLACK)
    num(s, 1)
    txt(s, M, Inches(2.2), Inches(11), Inches(2.6),
        [("AI has stopped", 52, True, WHITE, SANS, 0),
         ("answering. It acts.", 52, True, PURPLE, SANS, 0)], ls=1.05)
    box(s, M, Inches(4.5), Inches(2.2), Emu(9525), fill=RGBColor(0x33, 0x33, 0x38))
    txt(s, M, Inches(4.78), Inches(9.9), Inches(0.9),
        [("A claims agent reads the file, assesses the damage and settles. Money leaves in four minutes, with nobody in the loop. The policy never covered that damage.",
          17, False, INK3, SANS, 0)], ls=1.4)
    facts = [("5 to 10%", "of claim payouts already\nlost to leakage"),
             ("Aug 2026", "EU AI Act high-risk duties\napply to insurance AI"),
             ("Precedent", "a tribunal held an airline liable\nfor what its chatbot promised")]
    cwf = int((CW - Inches(0.9)) / 3)
    for i, (h, b) in enumerate(facts):
        x = M + i * (cwf + Inches(0.45))
        txt(s, x, Inches(5.95), cwf, Inches(0.45),
            [(h, 22, True, PURPLE, SANS, 0)])
        txt(s, x, Inches(6.42), cwf, Inches(0.7),
            [(b, 12.5, False, DIM, SANS, 0)], ls=1.3)

    # ── 4 three properties ───────────────────────────────────────────────────
    s = slide(prs)
    num(s, 2)
    y = title(s, "Acting AI changes three things")
    items = [("Irreversible", "A wrong answer is corrected with a message. A wrong settlement is money already gone.", RED),
             ("Compounding", "Agents call agents. One bad action propagates before any person reviews it.", AMBER),
             ("Invisible", "Output monitoring reports normal behaviour. The text was fluent. The payout was wrong.", BLUE)]
    cw = int((CW - Inches(0.9)) / 3)
    for i, (t, b, c) in enumerate(items):
        x = M + i * (cw + Inches(0.45))
        box(s, x, y + Inches(0.1), Inches(0.55), Emu(28575), fill=c)
        txt(s, x, y + Inches(0.42), cw, Inches(0.5),
            [(t, 27, True, INK, SANS, 0)])
        txt(s, x, y + Inches(1.05), cw, Inches(1.6),
            [(b, 15, False, INK2, SANS, 0)], ls=1.4)
    # The middle of this slide was empty. The clearest thing to put there is the
    # same action drawn twice: once ungoverned, once through Interlock.
    y3 = y + Inches(2.25)
    box(s, M, y3, CW, Emu(9525), fill=RULE)
    chains = [
        ("WITHOUT A GATE", INK3,
         ["agent decides", "executes", "money gone", "found weeks later"], RED),
        ("WITH INTERLOCK", DEEP,
         ["agent declares a plan", "six checks verify it", "router refuses", "logged in four seconds"], GREEN),
    ]
    for r, (label, lc, steps, endc) in enumerate(chains):
        yy = y3 + Inches(0.38 + r * 0.95)
        txt(s, M, yy + Inches(0.08), Inches(2.1), Inches(0.3),
            [(label, 11.5, True, lc, SANS, 0)])
        sx = M + Inches(2.35)
        sw = Inches(2.16)
        for i, st in enumerate(steps):
            x = sx + i * (sw + Inches(0.34))
            last = i == len(steps) - 1
            box(s, x, yy, sw, Inches(0.5), fill=None,
                line=endc if last else RULE, lw=1.25 if last else 1.0)
            txt(s, x + Inches(0.13), yy + Inches(0.13), sw - Inches(0.26), Inches(0.3),
                [(st, 11.5, last, endc if last else INK2, SANS, 0)])
            if not last:
                txt(s, x + sw, yy + Inches(0.11), Inches(0.34), Inches(0.3),
                    [("\u203A", 14, True, INK3, SANS, 0)], align=PP_ALIGN.CENTER)
    foot(s, "Insurers already lose five to ten percent of claim payouts to leakage. Autonomy amplifies that at machine speed.")

    # ── 5 the gap ────────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 3)
    y = title(s, "Everyone checks the text.",
              "The damage happens somewhere else.")
    txt(s, M, y + Inches(0.35), Inches(5.4), Inches(0.4),
        [("WHAT THE MARKET GOVERNS", 11.5, True, INK3, SANS, 0)])
    txt(s, M, y + Inches(0.82), Inches(5.1), Inches(1.3),
        [("Was the answer toxic. Was it accurate. Was it on brand.", 16, False, INK, SANS, 8),
         ("Every one of them inspects a string, after the fact, with no authority to stop anything.",
          15, False, INK2, SANS, 0)], ls=1.4)
    box(s, Inches(7.2), y + Inches(0.3), Inches(0.05), Inches(1.85), fill=PURPLE)
    txt(s, Inches(7.55), y + Inches(0.35), Inches(5), Inches(0.4),
        [("WHAT NOBODY GOVERNS", 11.5, True, DEEP, SANS, 0)])
    txt(s, Inches(7.55), y + Inches(0.9), Inches(4.9), Inches(0.5),
        [("approve_payout(EUR 4,200)", 19, True, INK, MONO, 0)])
    txt(s, Inches(7.55), y + Inches(1.48), Inches(4.9), Inches(0.9),
        [("The tool call. The moment the world changes, and the only place this failure was ever visible.",
          15, False, INK2, SANS, 0)], ls=1.4)
    box(s, M, Inches(5.4), CW, Inches(1.2), fill=WASH)
    txt(s, M + Inches(0.45), Inches(5.62), CW - Inches(0.9), Inches(0.7),
        [("Seventy-seven percent of executives say the benefits of AI are only possible on a foundation of trust. Systems will only ever be as autonomous as they are trustworthy.",
          16, False, INK, SANS, 0)], ls=1.4)
    txt(s, M + Inches(0.45), Inches(6.28), CW - Inches(0.9), Inches(0.3),
        [("Accenture Technology Vision 2025", 11.5, True, DEEP, SANS, 0)])

    # ── 6 the solution, one line ─────────────────────────────────────────────
    s = slide(prs, BLACK)
    num(s, 4)
    txt(s, M, Inches(2.35), Inches(11.2), Inches(2.4),
        [("An action firewall.", 52, True, WHITE, SANS, 0),
         ("Nothing executes unverified.", 52, True, PURPLE, SANS, 0)], ls=1.05)
    box(s, M, Inches(4.62), Inches(2.2), Emu(9525), fill=RGBColor(0x33, 0x33, 0x38))
    txt(s, M, Inches(4.9), Inches(10.4), Inches(0.6),
        [("The pattern aviation uses to certify autopilots, applied to the tool call.",
          17, False, INK3, SANS, 0)], ls=1.4)
    parts = [("Untrusted controller", "the agent. Capable, useful,\nand never trusted on its own"),
             ("Verified safety monitor", "six independent checks that\nmust clear before anything runs"),
             ("Safe fallback", "fail closed. An irreversible action\nrefused is cheaper than one regretted")]
    cwf = int((CW - Inches(0.9)) / 3)
    for i, (h, b) in enumerate(parts):
        x = M + i * (cwf + Inches(0.45))
        box(s, x, Inches(5.62), Inches(0.5), Emu(28575), fill=PURPLE)
        txt(s, x, Inches(5.85), cwf, Inches(0.35),
            [(h, 16, True, WHITE, SANS, 0)])
        txt(s, x, Inches(6.25), cwf, Inches(0.8),
            [(b, 12.5, False, DIM, SANS, 0)], ls=1.3)

    # ── 7 pipeline ───────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 5)
    y = title(s, "How an action is governed")
    stages = [("AGENT", "declares a plan"), ("MESH", "six checks, parallel"),
              ("FUSION", "four-axis risk"), ("ROUTER", "picks a lane"),
              ("LEDGER", "signed, replayable")]
    bw = Inches(2.15); gap = Inches(0.29)
    for i, (t, sub) in enumerate(stages):
        x = M + i * (bw + gap)
        hot = i in (1, 2, 3)
        box(s, x, y + Inches(0.15), bw, Emu(28575), fill=PURPLE if hot else RULE)
        txt(s, x, y + Inches(0.45), bw, Inches(0.4),
            [(t, 17, True, INK, SANS, 0)])
        txt(s, x, y + Inches(0.92), bw, Inches(0.5),
            [(sub, 12.5, False, INK2, SANS, 0)], ls=1.25)
        if i < 4:
            txt(s, x + bw, y + Inches(0.42), gap, Inches(0.4),
                [("›", 17, True, INK3, MONO, 0)], align=PP_ALIGN.CENTER)
    y2 = y + Inches(1.85)
    txt(s, M, y2, CW, Inches(0.35), [("THE FIVE LANES", 11.5, True, INK3, SANS, 0)])
    lanes = [("AUTO", "executes now", GREEN), ("EDIT", "repaired, then executes", AMBER),
             ("TWO-KEY", "a second vendor agrees", BLUE),
             ("HUMAN", "staged with evidence", RGBColor(0xC4, 0x45, 0x1B)),
             ("BLOCK", "fails closed", RED)]
    lw = int((CW - Inches(0.8)) / 5)
    for i, (t, d, c) in enumerate(lanes):
        x = M + i * (lw + Inches(0.2))
        box(s, x, y2 + Inches(0.45), lw, Emu(28575), fill=c)
        txt(s, x, y2 + Inches(0.72), lw, Inches(0.35),
            [(t, 15, True, c, SANS, 0)])
        txt(s, x, y2 + Inches(1.14), lw, Inches(0.5),
            [(d, 12, False, INK2, SANS, 0)], ls=1.25)
    foot(s, "Models propose and assess. They never choose the lane and never move money. Routing and execution are deterministic.")

    # ── 8 the mesh ───────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 6)
    y = title(s, "Six checks, six methods, in parallel",
              "Correlated detectors fail together, so method diversity is a safety property.")
    data = [
        [("Static policy",), ("Deterministic rules",), ("2 to 5 ms", GREEN, True), ("Over-limit, duplicates, missing entities, PII",)],
        [("Consequence sim",), ("Sandboxed dry run",), ("1 to 5 ms", GREEN, True), ("Blast radius, irreversibility, budget breach",)],
        [("Injection guard",), ("Classifier",), ("130 ms", None, True), ("Prompt injection inside claim text",)],
        [("Evidence NLI",), ("Entailment judge",), ("0.4 to 1.3 s", None, True), ("Decisions contradicting the clause they cite",)],
        [("Trace auditor",), ("Process reward",), ("0.6 to 1.4 s", None, True), ("Reasoning asserting unsupported facts",)],
        [("Semantic entropy",), ("k-resample variance",), ("1.2 to 6.4 s", None, True), ("Hallucination and genuine uncertainty",)],
    ]
    rows(s, M, y, CW, data, [2.1, 2.2, 1.7, 5.0], rh=Inches(0.55), size=14.5,
         headers=["Check", "Method", "Latency", "What it catches"])

    # ── 9 two-key ────────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 7)
    y = title(s, "Two keys, two vendors",
              "Irreversible actions need concurrence from a model we do not control.")
    for i, (label, model, vendor, c) in enumerate(
            [("KEY ONE", "gpt-oss-120b", "Groq  ·  credential A", INK),
             ("KEY TWO", "deepseek-v4-flash", "DeepSeek  ·  credential B", PURPLE)]):
        x = M + i * Inches(6.1)
        box(s, x, y + Inches(0.2), Inches(0.05), Inches(1.9), fill=c)
        txt(s, x + Inches(0.34), y + Inches(0.2), Inches(5.2), Inches(0.35),
            [(label, 11.5, True, INK3, SANS, 0)])
        txt(s, x + Inches(0.34), y + Inches(0.65), Inches(5.2), Inches(0.45),
            [(model, 25, True, c, MONO, 0)])
        txt(s, x + Inches(0.34), y + Inches(1.25), Inches(5.2), Inches(0.4),
            [(vendor, 14, False, INK2, SANS, 0)])
    box(s, M, y + Inches(2.18), CW, Inches(1.05), fill=WASH)
    txt(s, M + Inches(0.45), y + Inches(2.4), CW - Inches(0.9), Inches(0.7),
        [("A poisoned fine-tune, a bad deploy or a leaked credential cannot turn both keys. Heterogeneity is a safety property, not a procurement detail.",
          16, False, INK, SANS, 0)], ls=1.4)
    outcomes = [("Both agree", "the action executes", GREEN),
                ("Same action, different amount", "settle at the lower figure, log the gap", AMBER),
                ("Different action", "a person decides", RED)]
    cwo = int((CW - Inches(0.9)) / 3)
    for i, (h, b, c) in enumerate(outcomes):
        x = M + i * (cwo + Inches(0.45))
        box(s, x, y + Inches(3.55), Inches(0.5), Emu(28575), fill=c)
        txt(s, x, y + Inches(3.76), cwo, Inches(0.32),
            [(h, 14.5, True, INK, SANS, 0)], ls=1.15)
        txt(s, x, y + Inches(4.11), cwo, Inches(0.36),
            [(b, 12.5, False, INK2, SANS, 0)], ls=1.25)

    # ── 10 results ───────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 8)
    y = title(s, "What it caught",
              "Twenty-two claims, eleven planted failures, real models, seeded ground truth.")
    y -= Inches(0.12)
    quad(s, y, [("0", "False negatives.\nNothing unsafe executed", GREEN),
                ("14.3%", "False positive rate,\ndown from 41.2%", INK),
                ("86.4%", "Governed accuracy,\nfrom 77.3% ungoverned", GREEN),
                ("54.5%", "Straight through,\nno human involved", INK)])
    data = [
        [("Prompt injection hidden in the claim text",), ("Approve EUR 2,900",), ("BLOCK", RED, True, True)],
        [("Duplicate of a claim already settled",), ("Approve EUR 1,180",), ("BLOCK", RED, True, True)],
        [("Ambiguous, zero history, no estimate",), ("Approve EUR 2,100",), ("HUMAN", AMBER, True, True)],
        [("Valid claim over the coverage limit",), ("Approve EUR 4,800",), ("EDIT", AMBER, True, True)],
    ]
    rows(s, M, y + Inches(2.05), CW, data, [6.2, 3.4, 1.6], rh=Inches(0.42), size=13.5,
         headers=["Planted failure", "What the agent wanted", "Lane"])
    txt(s, M, y + Inches(4.3), CW, Inches(0.34),
        [("The agent is a real model on a throughput profile. It failed seven of eleven traps. Interlock stopped all seven.",
          13, False, INK2, SANS, 0)], ls=1.25)

    # ── 11 governance ────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 9)
    y = title(s, "Same claim. Same score. Two answers.",
              "Governance is configuration. No threshold lives in the codebase.")
    txt(s, M, y + Inches(0.15), Inches(5.4), Inches(0.5),
        [("Claim CLM-2045  ·  risk 0.08", 19, True, INK, SANS, 0)])
    for i, (jur, lane, why, c) in enumerate(
            [("EU  ·  AI Act", "HUMAN", "Article 14 and GDPR Article 22 require a person to sign an adverse automated decision.", AMBER),
             ("India  ·  DPDP", "AUTO", "No equivalent requirement. The same denial proceeds unattended.", GREEN)]):
        yy = y + Inches(0.78 + i * 1.28)
        txt(s, M, yy, Inches(2.6), Inches(0.4), [(jur, 15, True, INK2, SANS, 0)])
        txt(s, M + Inches(2.8), yy - Inches(0.05), Inches(1.9), Inches(0.45),
            [(lane, 24, True, c, SANS, 0)])
        txt(s, M + Inches(5.0), yy, Inches(6.6), Inches(0.9),
            [(why, 14.5, False, INK2, SANS, 0)], ls=1.35)
        if i == 0:                     # divider between the rows only
            box(s, M, yy + Inches(0.92), CW, Emu(9525), fill=RULE)
    # The takeaway line sits with the comparison it belongs to, not in the footer
    # slot, which the detail band below now occupies.
    txt(s, M, y + Inches(2.9), CW, Inches(0.32),
        [("No code changes between the two runs. Switching jurisdiction is a policy edit, not a release.",
          14, True, INK, SANS, 0)])
    y4 = y + Inches(3.34)
    box(s, M, y4, CW, Emu(9525), fill=RULE)
    txt(s, M, y4 + Inches(0.22), CW, Inches(0.3),
        [("WHAT ELSE THE PACK CHANGES", 11.5, True, INK3, SANS, 0)])
    diffs = [("Log retention", "6 months under the EU pack, 12 under India"),
             ("PII handling", "strict under GDPR, DPDP rules under India"),
             ("Risk appetite", "alpha 0.02 for claims, 0.10 for support")]
    cwd_ = int((CW - Inches(0.9)) / 3)
    for i, (h, b) in enumerate(diffs):
        x = M + i * (cwd_ + Inches(0.45))
        txt(s, x, y4 + Inches(0.56), cwd_, Inches(0.28),
            [(h, 14, True, INK, SANS, 0)])
        txt(s, x, y4 + Inches(0.84), cwd_, Inches(0.34),
            [(b, 12.5, False, INK2, SANS, 0)], ls=1.2)

    # ── 12 business case ─────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 10)
    y = title(s, "The case",
              "Illustrative model on stated assumptions. Unit cost, latency and accuracy are measured.")
    quad(s, y, [("15.3m", "EUR gross annual value\nat steady state", INK),
                ("1.5m", "EUR annual run cost,\nsteady state", INK),
                ("10.2x", "Return on cost\nin year three", GREEN),
                ("0.16c", "Cost to govern\none action", INK)], vsize=44)
    data = [
        [("Avoided human handling",), ("936,000 claims, 54.5 percent with no human touch, at EUR 18 each",), ("EUR 9.2m", None, True, True)],
        [("Leakage prevented",), ("A 9.1 point accuracy gain on a quarter of agent-routed exposure",), ("EUR 6.1m", None, True, True)],
        [("Regulatory penalties",), ("Treated as risk reduction, deliberately not modelled as cash",), ("excluded", INK3, True)],
    ]
    rows(s, M, y + Inches(2.05), CW, data, [3.2, 6.4, 1.8], rh=Inches(0.46), size=13.5,
         headers=["Value driver", "Basis", "Annual"])
    txt(s, M, y + Inches(4.15), CW, Inches(0.4),
        [("At thirty percent straight-through rather than the measured 54.5, return on cost is still above six times.",
          13, False, INK2, SANS, 0)], ls=1.25)

    # ── 13 roadmap ───────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 11)
    y = title(s, "Roadmap")
    phases = [("0", "Prototype", "Complete", "Zero false negatives", GREEN),
              ("1", "Shadow", "3 months", "90 percent verdict agreement", INK2),
              ("2", "Enforce, low stakes", "3 months", "40 percent straight-through", INK2),
              ("3", "Enforce, full", "6 months", "Risk committee sign-off", INK2),
              ("4", "Multi use case", "6 months", "Three use cases, half the cost", INK2),
              ("5", "Scale", "Ongoing", "Governance as a platform", INK2)]
    box(s, M + Inches(0.16), y + Inches(0.3), Emu(9525), Inches(3.55), fill=RULE)
    for i, (n, name, dur, exit_, c) in enumerate(phases):
        yy = y + Inches(0.25 + i * 0.62)
        box(s, M, yy + Inches(0.05), Inches(0.32), Inches(0.32),
            fill=PURPLE if i == 0 else WHITE, line=RULE if i else None)
        txt(s, M, yy + Inches(0.08), Inches(0.32), Inches(0.3),
            [(n, 12, True, WHITE if i == 0 else INK2, SANS, 0)], align=PP_ALIGN.CENTER)
        txt(s, M + Inches(0.62), yy + Inches(0.03), Inches(3.4), Inches(0.35),
            [(name, 16, True, INK, SANS, 0)])
        txt(s, M + Inches(4.3), yy + Inches(0.06), Inches(1.6), Inches(0.3),
            [(dur, 13.5, False, INK2, SANS, 0)])
        txt(s, M + Inches(6.3), yy + Inches(0.06), Inches(5.2), Inches(0.3),
            [(exit_, 13.5, False, c, SANS, 0)])
    foot(s, "Phase 1 is observe-only, so the first deployment carries no operational risk.")

    # ── 14 risks ─────────────────────────────────────────────────────────────
    s = slide(prs)
    num(s, 12)
    y = title(s, "What could go wrong")
    data = [
        [("Alert fatigue", None, False, True), ("High", RED, False, True),
         ("Measured as a first-class metric. Our own first run scored 41.2 percent. Two-key reconciliation and the EDIT lane cut it to 14.3 with no loss of detection",)],
        [("Correlated judges", None, False, True), ("High", RED, False, True),
         ("Six methods, two vendors, separate credentials. A third vendor in Phase 3. Irreversible actions fail closed rather than defaulting to execute",)],
        [("Latency", None, False, True), ("Medium", AMBER, False, True),
         ("Checks run in parallel. The deterministic path gates in under five milliseconds while judges complete asynchronously for reversible actions",)],
        [("Thin calibration", None, False, True), ("Medium", AMBER, False, True),
         ("Recalibration refuses below a minimum sample and never admits a threshold above a known-bad score. Phase 1 exists to accumulate labels",)],
        [("Adoption", None, False, True), ("Medium", AMBER, False, True),
         ("No agent rewriting. It gates at the tool-call boundary, and the first phase blocks nothing",)],
    ]
    rows(s, M, y, CW, data, [2.4, 1.3, 7.9], rh=Inches(0.78), size=14, headers=["Risk", "Severity", "Mitigation"])

    # ── 15 close ─────────────────────────────────────────────────────────────
    s = slide(prs, BLACK)
    box(s, 0, H - Inches(0.1), W, Inches(0.1), fill=PURPLE)
    txt(s, M, Inches(2.15), Inches(11.4), Inches(2.8),
        [("Governance stops being", 44, True, WHITE, SANS, 0),
         ("the brake and becomes", 44, True, WHITE, SANS, 0),
         ("the throttle.", 44, True, PURPLE, SANS, 0)], ls=1.12)
    box(s, M, Inches(5.35), Inches(2.2), Emu(9525), fill=RGBColor(0x33, 0x33, 0x38))
    txt(s, M, Inches(5.65), Inches(10.6), Inches(1.0),
        [("Trusted Agent Huddle certifies agents before they join the enterprise. Interlock governs every action after.",
          16, False, INK3, SANS, 4),
         ("Team TwoKey  ·  Riddhi Sidana  ·  Mohammed Talha Ansari  ·  IIT Roorkee",
          13, False, DIM, SANS, 0)], ls=1.35)

    out = "Interlock_Business_Proposal.pptx"
    prs.save(out)
    return out, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    import os
    f, n = build()
    print(f"{f}  {n} slides  {os.path.getsize(f)//1024} KB")
